from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PptPt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from io import BytesIO
import re
import requests

# ... (parse_markdown_to_docx and export_to_docx remain same)

def export_to_pptx(project_data: dict) -> BytesIO:
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project_data['title']
    subtitle.text = project_data.get('topic', '')

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5] # Good for charts/images
    
    for item in sorted(project_data['items'], key=lambda x: x['order']):
        if item.get('type') == 'chart':
            slide = prs.slides.add_slide(title_only_layout)
            title = slide.shapes.title
            title.text = item['title']
            
            chart_data_json = item.get('chart_data')
            if chart_data_json:
                chart_data = CategoryChartData()
                chart_data.categories = chart_data_json.get('categories', [])
                for series in chart_data_json.get('series', []):
                    chart_data.add_series(series['name'], series['values'])
                
                x, y, cx, cy = PptPt(50), PptPt(100), PptPt(600), PptPt(400)
                chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                if chart_data_json.get('type') == 'pie':
                    chart_type = XL_CHART_TYPE.PIE
                elif chart_data_json.get('type') == 'line':
                    chart_type = XL_CHART_TYPE.LINE
                
                slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
                
        elif item.get('type') == 'image_prompt':
            slide = prs.slides.add_slide(title_only_layout)
            title = slide.shapes.title
            title.text = item['title']
            
            image_url = item.get('image_url')
            if image_url:
                try:
                    response = requests.get(image_url)
                    image_stream = BytesIO(response.content)
                    
                    left = PptPt(100)
                    top = PptPt(100)
                    width = PptPt(500)
                    # height will be auto-calculated to preserve aspect ratio if not specified, 
                    # but let's set a max height or width
                    
                    slide.shapes.add_picture(image_stream, left, top, width=width)
                except Exception as e:
                    print(f"Error downloading image: {e}")
                    # Fallback to text placeholder
                    left = top = width = height = PptPt(100)
                    txBox = slide.shapes.add_textbox(left, top, PptPt(500), PptPt(300))
                    tf = txBox.text_frame
                    tf.text = f"[IMAGE DOWNLOAD FAILED]\n\nPrompt: {item.get('image_prompt')}"
            else:
                # Fallback to placeholder if no URL
                left = top = width = height = PptPt(100)
                txBox = slide.shapes.add_textbox(left, top, PptPt(500), PptPt(300))
                tf = txBox.text_frame
                tf.text = f"[IMAGE PLACEHOLDER]\n\nPrompt: {item.get('image_prompt')}"
            
        else:
            # Standard Slide (Text + Optional Image)
            image_url = item.get('image_url')
            
            if image_url:
                # Use Two Content Layout (Index 3)
                two_content_layout = prs.slide_layouts[3] 
                slide = prs.slides.add_slide(two_content_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = item['title']
                
                # Left side: Text
                left_body = shapes.placeholders[1]
                tf = left_body.text_frame
                tf.clear()
                parse_markdown_to_pptx(tf, item['content'])
                
                # Right side: Image
                right_body = shapes.placeholders[2]
                try:
                    response = requests.get(image_url)
                    image_stream = BytesIO(response.content)
                    right_body.insert_picture(image_stream)
                except Exception as e:
                    print(f"Error downloading image for slide: {e}")
                    right_body.text = f"[IMAGE FAILED]\n{item.get('image_prompt', '')}"
            else:
                # Standard Bullet Layout
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = item['title']
                tf = body_shape.text_frame
                tf.clear()
                parse_markdown_to_pptx(tf, item['content'])

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def parse_markdown_to_docx(doc, text):
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('- ') or line.startswith('* '):
            p = doc.add_paragraph(line[2:], style='List Bullet')
        else:
            p = doc.add_paragraph()
            # Simple bold parsing **text**
            parts = re.split(r'(\*\*.*?\*\*)', line)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    p.add_run(part)

def export_to_docx(project_data: dict) -> BytesIO:
    doc = Document()
    doc.add_heading(project_data['title'], 0)
    
    for item in sorted(project_data['items'], key=lambda x: x['order']):
        doc.add_heading(item['title'], level=1)
        parse_markdown_to_docx(doc, item['content'])
        
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def parse_markdown_to_pptx(text_frame, text):
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        p = text_frame.add_paragraph()
        
        # Handle bullets
        clean_line = line
        if line.startswith('- ') or line.startswith('* '):
            clean_line = line[2:]
            p.level = 0 # Top level bullet
        elif line.startswith('  - ') or line.startswith('  * '):
            clean_line = line[4:]
            p.level = 1
            
        # Simple bold parsing
        parts = re.split(r'(\*\*.*?\*\*)', clean_line)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = p.add_run()
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run = p.add_run()
                run.text = part

def export_to_pptx(project_data: dict) -> BytesIO:
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project_data['title']
    subtitle.text = project_data.get('topic', '')

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for item in sorted(project_data['items'], key=lambda x: x['order']):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = item['title']
        tf = body_shape.text_frame
        tf.clear() # Clear default placeholder text
        
        parse_markdown_to_pptx(tf, item['content'])

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
